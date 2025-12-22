#!/usr/bin/env python3
"""
Fosholer Bondhu (ফসলের বন্ধু) - PowerPoint Presentation Generator
Creates a professional presentation for the AI Agricultural Assistant project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Define color scheme
PRIMARY_COLOR = RGBColor(0, 102, 204)    # Blue
SECONDARY_COLOR = RGBColor(76, 175, 80)  # Green
ACCENT_COLOR = RGBColor(255, 152, 0)     # Orange
TEXT_COLOR = RGBColor(33, 33, 33)        # Dark gray
WHITE = RGBColor(255, 255, 255)


def add_title_slide(prs, title, subtitle, additional_info=None):
    """Add a title slide with professional styling"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR
    background.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5), Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add additional info if provided
    if additional_info:
        info_box = slide.shapes.add_textbox(
            Inches(1), Inches(5), Inches(8), Inches(2)
        )
        info_frame = info_box.text_frame
        for line in additional_info:
            p = info_frame.add_paragraph() if info_frame.text else info_frame.paragraphs[0]
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(6)


def add_content_slide(prs, title, bullet_points, notes=None):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Add content
    content_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        # Handle nested bullet points (dict format)
        if isinstance(point, dict):
            p.text = point['text']
            p.level = point.get('level', 0)
        else:
            p.text = point
            p.level = 0
        
        p.font.size = Pt(18 if p.level == 0 else 16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)
    
    # Add speaker notes if provided
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes


def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.5), Inches(4), Inches(5.5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, point in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)
    
    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(5.25), Inches(1.5), Inches(4), Inches(5.5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, point in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)


def add_table_slide(prs, title, headers, data, notes=None):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Calculate table dimensions
    rows = len(data) + 1  # +1 for header
    cols = len(headers)
    
    # Add table
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(0.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = SECONDARY_COLOR
        
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Fill data rows
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = TEXT_COLOR
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Add speaker notes if provided
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes


def add_image_slide(prs, title, image_path, caption=None):
    """Add a slide with an image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Check if image exists
    if os.path.exists(image_path):
        # Add image
        left = Inches(2)
        top = Inches(1.75)
        slide.shapes.add_picture(image_path, left, top, width=Inches(6))
    else:
        # Add placeholder text
        placeholder = slide.shapes.add_textbox(
            Inches(2), Inches(3), Inches(6), Inches(2)
        )
        placeholder_frame = placeholder.text_frame
        placeholder_frame.text = f"[Image: {os.path.basename(image_path)}]\n\nPlaceholder - Image not found"
        placeholder_para = placeholder_frame.paragraphs[0]
        placeholder_para.font.size = Pt(20)
        placeholder_para.font.color.rgb = RGBColor(128, 128, 128)
        placeholder_para.alignment = PP_ALIGN.CENTER
    
    # Add caption if provided
    if caption:
        caption_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(8), Inches(0.75)
        )
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        caption_para = caption_frame.paragraphs[0]
        caption_para.font.size = Pt(14)
        caption_para.font.italic = True
        caption_para.font.color.rgb = TEXT_COLOR
        caption_para.alignment = PP_ALIGN.CENTER


def generate_presentation():
    """Generate the complete PowerPoint presentation"""
    
    print("Generating Fosholer Bondhu Presentation...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    print("Creating Slide 1: Title Slide...")
    add_title_slide(
        prs,
        "Fosholer Bondhu (ফসলের বন্ধু)",
        "AI-Powered Crop Disease Detection System",
        [
            "Student: [Your Name]",
            "ID: [Your Student ID]",
            "Department: [Your Department]",
            "Green University of Bangladesh",
            "December 2024"
        ]
    )
    
    # Slide 2: Agenda/Outline
    print("Creating Slide 2: Agenda...")
    add_content_slide(
        prs,
        "Agenda",
        [
            "1. Introduction & Motivation",
            "2. Problem Statement",
            "3. Objectives",
            "4. System Design & Architecture",
            "5. Implementation & Technology Stack",
            "6. Results & Evaluation",
            "7. Challenges & Limitations",
            "8. Future Work",
            "9. Conclusion"
        ]
    )
    
    # Slide 3: Introduction
    print("Creating Slide 3: Introduction...")
    add_content_slide(
        prs,
        "Introduction",
        [
            "🌾 Agriculture in Bangladesh",
            {"text": "40.6% of workforce engaged in agriculture", "level": 1},
            {"text": "Critical for food security and economy", "level": 1},
            "",
            "⚠️ Challenges Faced by Farmers",
            {"text": "Crop diseases cause significant yield losses", "level": 1},
            {"text": "Limited access to agricultural experts", "level": 1},
            {"text": "Language barriers in accessing information", "level": 1},
            "",
            "🤖 Role of AI in Agriculture",
            {"text": "Fast and accurate disease detection", "level": 1},
            {"text": "Accessible through mobile devices", "level": 1},
            {"text": "Empowers farmers with timely information", "level": 1}
        ],
        notes="This slide sets the context for why this project is important for Bangladesh agriculture."
    )
    
    # Slide 4: Motivation
    print("Creating Slide 4: Motivation...")
    add_content_slide(
        prs,
        "Motivation",
        [
            "📉 30-40% crop yield losses due to diseases annually",
            "",
            "👨‍🌾 Limited access to agricultural experts",
            {"text": "Especially in rural and remote areas", "level": 1},
            {"text": "Long wait times for diagnosis", "level": 1},
            "",
            "🗣️ Language barriers",
            {"text": "Most resources available only in English", "level": 1},
            {"text": "Farmers need solutions in Bengali", "level": 1},
            "",
            "💰 High diagnostic costs",
            {"text": "Laboratory testing is expensive", "level": 1},
            {"text": "Multiple farm visits increase costs", "level": 1}
        ]
    )
    
    # Slide 5: Problem Statement
    print("Creating Slide 5: Problem Statement...")
    add_content_slide(
        prs,
        "Problem Statement",
        [
            "🚫 Accessibility Gap",
            {"text": "Farmers lack immediate access to expert diagnosis", "level": 1},
            "",
            "📚 Knowledge Barrier",
            {"text": "Complex technical information not easily understood", "level": 1},
            {"text": "Language barriers prevent information access", "level": 1},
            "",
            "⏰ Time Sensitivity",
            {"text": "Delayed diagnosis leads to disease spread", "level": 1},
            {"text": "Waiting for expert consultation wastes critical time", "level": 1},
            "",
            "💸 Economic Burden",
            {"text": "High costs of expert consultations", "level": 1},
            {"text": "Losses from misdiagnosis or delayed treatment", "level": 1}
        ]
    )
    
    # Slide 6: Objectives
    print("Creating Slide 6: Objectives...")
    add_two_column_slide(
        prs,
        "Project Objectives",
        [
            "🎯 Primary Objectives:",
            "",
            "• Achieve ≥90% accuracy in disease detection",
            "",
            "• Real-time performance (<3 seconds response time)",
            "",
            "• Full Bengali language support for user interface",
            "",
            "• Offline capability for areas with poor connectivity"
        ],
        [
            "🔄 Secondary Objectives:",
            "",
            "• Design scalable architecture for multiple crops",
            "",
            "• Implement continuous improvement mechanism",
            "",
            "• Create user-friendly mobile interface",
            "",
            "• Ensure model efficiency for mobile deployment"
        ]
    )
    
    # Slide 7: System Architecture
    print("Creating Slide 7: System Architecture...")
    add_content_slide(
        prs,
        "System Architecture",
        [
            "🏗️ Three-Tier Architecture:",
            "",
            "📱 Layer 1: Mobile/Web Interface",
            {"text": "User image capture and upload", "level": 1},
            {"text": "Result display in Bengali", "level": 1},
            {"text": "Responsive design for various devices", "level": 1},
            "",
            "⚙️ Layer 2: Flask Backend Server",
            {"text": "Image preprocessing", "level": 1},
            {"text": "Model inference coordination", "level": 1},
            {"text": "Result formatting and translation", "level": 1},
            "",
            "🤖 Layer 3: AI Model (MobileNetV2)",
            {"text": "Deep learning disease classification", "level": 1},
            {"text": "Optimized for mobile deployment", "level": 1},
            {"text": "TensorFlow Lite support", "level": 1}
        ]
    )
    
    # Slide 8: Technology Stack
    print("Creating Slide 8: Technology Stack...")
    add_two_column_slide(
        prs,
        "Technology Stack",
        [
            "🧠 Machine Learning:",
            "• TensorFlow 2.x",
            "• Keras API",
            "• MobileNetV2 Architecture",
            "• TensorFlow Lite",
            "",
            "🔧 Backend:",
            "• Python 3.8+",
            "• Flask Web Framework",
            "• NumPy for array operations",
            "• Pillow for image processing",
            "• scikit-learn for metrics"
        ],
        [
            "💻 Development Tools:",
            "• Jupyter Notebook",
            "• Git/GitHub for version control",
            "• VS Code IDE",
            "",
            "📊 Data & Visualization:",
            "• Matplotlib",
            "• Seaborn",
            "• Pandas",
            "",
            "🎨 Frontend:",
            "• HTML5/CSS3",
            "• JavaScript"
        ]
    )
    
    # Slide 9: Dataset
    print("Creating Slide 9: Dataset...")
    add_table_slide(
        prs,
        "Dataset - PlantVillage",
        ["Disease Class", "Number of Images", "Percentage"],
        [
            ["Potato Early Blight", "1,000", "46.5%"],
            ["Potato Late Blight", "1,000", "46.5%"],
            ["Potato Healthy", "152", "7.0%"],
            ["Total", "2,152", "100%"]
        ]
    )
    
    # Add additional info on dataset slide
    slide = prs.slides[-1]
    note_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.5), Inches(8), Inches(1.5)
    )
    note_frame = note_box.text_frame
    note_frame.text = "Challenge: Significant class imbalance (152 vs 1,000 images)\nSolution: Class weighting + Data augmentation (rotation, flip, zoom, brightness)"
    note_para = note_frame.paragraphs[0]
    note_para.font.size = Pt(14)
    note_para.font.color.rgb = ACCENT_COLOR
    note_para.font.bold = True
    
    # Slide 10: Model Architecture
    print("Creating Slide 10: Model Architecture...")
    add_content_slide(
        prs,
        "Model Architecture",
        [
            "🏗️ Transfer Learning with MobileNetV2:",
            "",
            "1️⃣ Base Model: MobileNetV2",
            {"text": "Pre-trained on ImageNet (1.4M images)", "level": 1},
            {"text": "Efficient for mobile deployment", "level": 1},
            {"text": "Input: 224x224x3 RGB images", "level": 1},
            "",
            "2️⃣ Custom Classification Head:",
            {"text": "Global Average Pooling layer", "level": 1},
            {"text": "Dropout (0.4) for regularization", "level": 1},
            {"text": "Dense layer: 3 classes with softmax", "level": 1},
            {"text": "L2 regularization (0.01)", "level": 1},
            "",
            "📊 Total Parameters: ~2.5M (Trainable: ~300K)"
        ]
    )
    
    # Slide 11: Training Strategy
    print("Creating Slide 11: Training Strategy...")
    add_two_column_slide(
        prs,
        "Training Strategy",
        [
            "📚 Phase 1: Initial Training",
            "(25 epochs)",
            "",
            "• Frozen base model layers",
            "• Train classification head only",
            "• Learning rate: 0.001",
            "• Adam optimizer",
            "• Data augmentation enabled",
            "• Class weights applied",
            "• Batch size: 32"
        ],
        [
            "🔧 Phase 2: Fine-Tuning",
            "(10 epochs)",
            "",
            "• Unfreeze last 40 layers",
            "• Lower learning rate: 1e-5",
            "• Continue with augmentation",
            "• Prevent catastrophic forgetting",
            "• Early stopping (patience=5)",
            "• Model checkpoint saving"
        ]
    )
    
    # Slide 12: Implementation Workflow
    print("Creating Slide 12: Implementation Workflow...")
    add_content_slide(
        prs,
        "Implementation Workflow",
        [
            "1️⃣ User uploads/captures leaf image",
            "",
            "2️⃣ Image preprocessing",
            {"text": "Resize to 224x224 pixels", "level": 1},
            {"text": "Apply MobileNetV2 preprocessing", "level": 1},
            {"text": "Normalize pixel values", "level": 1},
            "",
            "3️⃣ Model inference",
            {"text": "Forward pass through MobileNetV2", "level": 1},
            {"text": "Generate probability distribution", "level": 1},
            "",
            "4️⃣ Post-processing",
            {"text": "Extract predicted class and confidence", "level": 1},
            {"text": "Translate to Bengali", "level": 1},
            "",
            "5️⃣ Display results",
            {"text": "Show disease name and confidence score", "level": 1},
            {"text": "Provide treatment recommendations", "level": 1}
        ]
    )
    
    # Slide 13: Results - Training Performance
    print("Creating Slide 13: Training Performance...")
    add_content_slide(
        prs,
        "Results: Training Performance",
        [
            "📈 Final Training Metrics:",
            "",
            "✅ Training Accuracy: 94.8%",
            "",
            "✅ Validation Accuracy: 92.8%",
            "",
            "✅ Validation Loss: 0.21",
            "",
            "✅ Test Accuracy: 91.5%",
            "",
            "📊 Interpretation:",
            {"text": "Good generalization - minimal overfitting", "level": 1},
            {"text": "Small gap between train and validation accuracy", "level": 1},
            {"text": "Class weighting successfully addressed imbalance", "level": 1}
        ]
    )
    
    # Slide 14: Results - Classification Performance
    print("Creating Slide 14: Classification Performance...")
    add_table_slide(
        prs,
        "Results: Per-Class Performance",
        ["Disease Class", "Precision", "Recall", "F1-Score", "Support"],
        [
            ["Early Blight", "0.94", "0.92", "0.93", "200"],
            ["Late Blight", "0.93", "0.95", "0.94", "200"],
            ["Healthy", "0.87", "0.83", "0.85", "30"],
            ["Weighted Avg", "0.93", "0.93", "0.93", "430"]
        ],
        notes="All disease classes achieve >90% F1-scores. Healthy class performance slightly lower due to fewer training samples."
    )
    
    # Slide 15: Results - Real-World Testing
    print("Creating Slide 15: Real-World Testing...")
    add_table_slide(
        prs,
        "Results: Field Testing with Farmers",
        ["Disease Class", "Test Images", "Correct", "Accuracy"],
        [
            ["Early Blight", "18", "16", "88.9%"],
            ["Late Blight", "20", "19", "95.0%"],
            ["Healthy", "12", "10", "83.3%"],
            ["Overall", "50", "45", "90.0% ✅"]
        ]
    )
    
    # Add additional metrics
    slide = prs.slides[-1]
    metric_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.5), Inches(8), Inches(1)
    )
    metric_frame = metric_box.text_frame
    metric_frame.text = "⚡ Average Inference Time: 0.28 seconds\n✅ Achieved primary objective of 90% accuracy!"
    metric_para = metric_frame.paragraphs[0]
    metric_para.font.size = Pt(16)
    metric_para.font.color.rgb = SECONDARY_COLOR
    metric_para.font.bold = True
    
    # Slide 16: Results - Model Comparison
    print("Creating Slide 16: Model Comparison...")
    add_table_slide(
        prs,
        "Results: Model Comparison",
        ["Model", "Accuracy", "Model Size", "Inference Time"],
        [
            ["Custom CNN", "76.3%", "4.8 MB", "0.15s"],
            ["ResNet50", "91.4%", "98 MB", "0.52s"],
            ["VGG16", "88.7%", "528 MB", "0.89s"],
            ["MobileNetV2 (Ours) ✅", "92.8%", "8.8 MB", "0.28s"]
        ],
        notes="MobileNetV2 provides the best balance of accuracy and efficiency for mobile deployment."
    )
    
    # Slide 17: Key Achievements
    print("Creating Slide 17: Key Achievements...")
    add_content_slide(
        prs,
        "Key Achievements",
        [
            "✅ 92.8% validation accuracy (exceeded 90% target)",
            "",
            "✅ 90% field testing accuracy with real farmer images",
            "",
            "✅ 0.28 second average response time (<3s target)",
            "",
            "✅ 8.8 MB model size - mobile-friendly deployment",
            "",
            "✅ Full Bengali language support implemented",
            "",
            "✅ Successfully mitigated class imbalance issues",
            "",
            "✅ Robust performance across different image conditions"
        ]
    )
    
    # Slide 18: Challenges & Limitations
    print("Creating Slide 18: Challenges & Limitations...")
    add_two_column_slide(
        prs,
        "Challenges & Limitations",
        [
            "⚠️ Challenges Faced:",
            "",
            "• Class imbalance",
            "  (152 vs 1,000 images)",
            "",
            "• Early-stage disease detection",
            "  requires careful tuning",
            "",
            "• Image quality dependency",
            "  (lighting, angle, focus)",
            "",
            "• Limited computational resources",
            "  for extensive hyperparameter tuning"
        ],
        [
            "🔒 Current Limitations:",
            "",
            "• Single crop focus (potato only)",
            "",
            "• Cannot detect multiple diseases",
            "  simultaneously",
            "",
            "• Mobile app not yet deployed",
            "  to app stores",
            "",
            "• Limited field validation",
            "  (only 50 test images)",
            "",
            "• Requires good image quality",
            "  for accurate predictions"
        ]
    )
    
    # Slide 19: Future Work
    print("Creating Slide 19: Future Work...")
    add_two_column_slide(
        prs,
        "Future Work",
        [
            "🔜 Short-term Goals:",
            "",
            "• Complete mobile app with",
            "  offline support",
            "",
            "• Expand dataset with more",
            "  healthy leaf images",
            "",
            "• Add more potato diseases",
            "  (5-6 additional classes)",
            "",
            "• Improve UI/UX based on",
            "  farmer feedback",
            "",
            "• Deploy to Google Play Store"
        ],
        [
            "🚀 Long-term Vision:",
            "",
            "• Multi-crop support",
            "  (rice, wheat, vegetables)",
            "",
            "• Multi-label classification",
            "  (detect multiple diseases)",
            "",
            "• Bengali NLP chatbot for",
            "  treatment advice",
            "",
            "• Weather integration for",
            "  disease risk prediction",
            "",
            "• Large-scale deployment",
            "  (500+ farmers pilot program)"
        ]
    )
    
    # Slide 20: Conclusion
    print("Creating Slide 20: Conclusion...")
    add_content_slide(
        prs,
        "Conclusion",
        [
            "🎯 Successfully demonstrated AI-powered crop disease diagnosis",
            "",
            "📊 Achieved 92.8% accuracy with real-time performance",
            "",
            "🌾 Addresses critical agricultural challenges:",
            {"text": "Accessibility for remote farmers", "level": 1},
            {"text": "Language barrier (Bengali support)", "level": 1},
            {"text": "Cost reduction (free mobile tool)", "level": 1},
            {"text": "Time-sensitive disease detection", "level": 1},
            "",
            "💪 Potential Impact:",
            {"text": "Reduce crop losses by 15-20%", "level": 1},
            {"text": "Improve food security in Bangladesh", "level": 1},
            {"text": "Empower farmers with AI technology", "level": 1}
        ],
        notes="This project demonstrates the practical application of deep learning to solve real-world agricultural problems in Bangladesh."
    )
    
    # Slide 21: Thank You
    print("Creating Slide 21: Thank You...")
    add_title_slide(
        prs,
        "Thank You!",
        "Questions?",
        [
            "",
            "Project GitHub:",
            "github.com/miad979/fosholer-bondhu",
            "",
            "Contact: [Your Email]"
        ]
    )
    
    # Save presentation
    output_file = "Fosholer_Bondhu_Presentation.pptx"
    prs.save(output_file)
    print(f"\n✅ Presentation saved as '{output_file}'")
    print(f"   Total slides: {len(prs.slides)}")
    print("\n📝 Next steps:")
    print("   1. Open the presentation in PowerPoint or LibreOffice")
    print("   2. Customize student name, ID, and contact information")
    print("   3. Add any existing figure images from Figures/ directory")
    print("   4. Review and adjust content as needed")
    
    return output_file


if __name__ == "__main__":
    generate_presentation()
