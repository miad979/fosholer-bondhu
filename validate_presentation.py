#!/usr/bin/env python3
"""
Presentation Validator
Validates the generated PowerPoint presentation and provides a detailed report.
"""

from pptx import Presentation
import sys
import os


def validate_presentation(filename='Fosholer_Bondhu_Presentation.pptx'):
    """Validate the presentation file"""
    
    if not os.path.exists(filename):
        print(f"❌ Error: File '{filename}' not found!")
        print(f"   Run 'python generate_presentation.py' first.")
        return False
    
    try:
        prs = Presentation(filename)
    except Exception as e:
        print(f"❌ Error loading presentation: {e}")
        return False
    
    print("=" * 70)
    print("📊 FOSHOLER BONDHU PRESENTATION VALIDATION REPORT")
    print("=" * 70)
    print()
    
    # Basic info
    print("📁 File Information:")
    print(f"   Filename: {filename}")
    print(f"   File size: {os.path.getsize(filename) / 1024:.1f} KB")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Dimensions: {prs.slide_width.inches}\" × {prs.slide_height.inches}\"")
    print()
    
    # Slide-by-slide analysis
    print("📋 Slide Contents:")
    print("-" * 70)
    
    expected_slides = [
        "Title Slide (with project name and Bengali text)",
        "Agenda/Outline",
        "Introduction",
        "Motivation",
        "Problem Statement",
        "Objectives (Primary & Secondary)",
        "System Architecture",
        "Technology Stack",
        "Dataset Information",
        "Model Architecture",
        "Training Strategy",
        "Implementation Workflow",
        "Results: Training Performance",
        "Results: Classification Performance",
        "Results: Real-World Testing",
        "Results: Model Comparison",
        "Key Achievements",
        "Challenges & Limitations",
        "Future Work",
        "Conclusion",
        "Thank You Slide"
    ]
    
    issues = []
    
    for i, slide in enumerate(prs.slides, 1):
        # Extract slide title and content
        title = ""
        text_count = 0
        has_table = False
        shape_types = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame.text:
                text = shape.text_frame.text.strip()
                if text:
                    text_count += 1
                    if not title and len(text) < 100:
                        title = text
            
            if shape.shape_type == 19:  # Table
                has_table = True
            
            shape_types.append(shape.shape_type)
        
        # Format output
        if not title:
            title = f"[Slide {i}]"
        
        status = "✅"
        details = []
        
        if text_count == 0:
            status = "⚠️"
            details.append("No text found")
            issues.append(f"Slide {i}: No text content")
        
        if has_table:
            details.append("Contains table")
        
        print(f"{status} Slide {i:2d}: {title[:50]}")
        if details:
            print(f"           {', '.join(details)}")
    
    print()
    print("-" * 70)
    
    # Summary
    print()
    print("📊 Summary:")
    print(f"   Expected slides: {len(expected_slides)}")
    print(f"   Actual slides: {len(prs.slides)}")
    
    if len(prs.slides) == 21:
        print("   ✅ Slide count matches expected (21 slides)")
    else:
        print(f"   ⚠️  Slide count mismatch")
        issues.append(f"Expected 21 slides, found {len(prs.slides)}")
    
    # Check for required content
    print()
    print("🔍 Content Validation:")
    
    checks = [
        ("Title slide with Bengali text", "ফসলের বন্ধু", True),
        ("Dataset information", "PlantVillage", False),
        ("Model architecture", "MobileNetV2", False),
        ("Training metrics", "92.8%", False),
        ("Field testing results", "90%", False),
        ("GitHub repository", "github.com/miad979", False),
    ]
    
    for check_name, search_term, required in checks:
        found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    if search_term.lower() in shape.text_frame.text.lower():
                        found = True
                        break
            if found:
                break
        
        if found:
            print(f"   ✅ {check_name}")
        else:
            status = "❌" if required else "⚠️"
            print(f"   {status} {check_name} - not found")
            if required:
                issues.append(f"Required content missing: {check_name}")
    
    print()
    
    # Final verdict
    if not issues:
        print("=" * 70)
        print("✅ VALIDATION PASSED - Presentation is ready!")
        print("=" * 70)
        print()
        print("📝 Next Steps:")
        print("   1. Open the presentation in PowerPoint/LibreOffice")
        print("   2. Customize student name, ID, and contact information")
        print("   3. Review content and adjust as needed")
        print("   4. Add figure images if available")
        print("   5. Practice your presentation!")
        print()
        return True
    else:
        print("=" * 70)
        print("⚠️  VALIDATION COMPLETED WITH WARNINGS")
        print("=" * 70)
        print()
        print("Issues found:")
        for issue in issues:
            print(f"   • {issue}")
        print()
        print("The presentation was generated but may need adjustments.")
        print()
        return True


if __name__ == "__main__":
    filename = sys.argv[1] if len(sys.argv) > 1 else 'Fosholer_Bondhu_Presentation.pptx'
    success = validate_presentation(filename)
    sys.exit(0 if success else 1)
